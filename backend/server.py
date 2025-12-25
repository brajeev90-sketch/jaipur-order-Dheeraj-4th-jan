from fastapi import FastAPI, APIRouter, HTTPException, UploadFile, File
from fastapi.responses import StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional
import uuid
from datetime import datetime, timezone
import io
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.pdfgen import canvas
from reportlab.lib.colors import HexColor
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, Table, TableStyle
from pptx import Presentation
from pptx.util import Inches, Pt
import re
import base64
import pandas as pd
import httpx

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Create the main app
app = FastAPI()

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")

# ============ MODELS ============

class OrderItem(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    product_code: str = ""
    description: str = ""
    category: str = ""
    height_cm: float = 0
    depth_cm: float = 0
    width_cm: float = 0
    cbm: float = 0
    cbm_auto: bool = True
    quantity: int = 1
    in_house_production: bool = True
    machine_hall: str = ""
    leather_code: str = ""
    leather_image: str = ""
    finish_code: str = ""
    finish_image: str = ""
    color_notes: str = ""
    leg_color: str = ""
    wood_finish: str = ""
    notes: str = ""
    images: List[str] = []
    reference_images: List[str] = []

class Order(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    sales_order_ref: str = ""
    buyer_po_ref: str = ""
    buyer_name: str = ""
    entry_date: str = ""
    factory_inform_date: str = ""
    status: str = "Draft"
    factory: str = ""
    items: List[OrderItem] = []
    created_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())
    updated_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())

class OrderCreate(BaseModel):
    sales_order_ref: str = ""
    buyer_po_ref: str = ""
    buyer_name: str = ""
    entry_date: str = ""
    factory_inform_date: str = ""
    status: str = "Draft"
    factory: str = ""
    items: List[OrderItem] = []

class OrderUpdate(BaseModel):
    sales_order_ref: Optional[str] = None
    buyer_po_ref: Optional[str] = None
    buyer_name: Optional[str] = None
    entry_date: Optional[str] = None
    factory_inform_date: Optional[str] = None
    status: Optional[str] = None
    factory: Optional[str] = None
    items: Optional[List[OrderItem]] = None

class LeatherLibraryItem(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    code: str
    name: str
    description: str = ""
    color: str = ""
    image: str = ""
    created_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())

class FinishLibraryItem(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    code: str
    name: str
    description: str = ""
    color: str = ""
    image: str = ""
    created_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())

class TemplateSettings(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = "default"
    company_name: str = "JAIPUR – A fine wood furniture company"
    logo_text: str = "JAIPUR"
    primary_color: str = "#3d2c1e"
    accent_color: str = "#d4622e"
    font_family: str = "Playfair Display, serif"
    body_font: str = "Manrope, sans-serif"
    page_margin_mm: int = 15
    show_borders: bool = True
    header_height_mm: int = 25
    footer_height_mm: int = 20

class ExportRecord(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    order_id: str
    export_type: str
    filename: str
    created_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())

class Product(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    product_code: str
    description: str = ""
    category: str = ""
    size: str = ""
    height_cm: float = 0
    depth_cm: float = 0
    width_cm: float = 0
    cbm: float = 0
    fob_price_usd: float = 0
    fob_price_gbp: float = 0
    warehouse_price_1: float = 0
    warehouse_price_2: float = 0
    image: str = ""
    images: List[str] = []
    created_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())
    updated_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())

class ProductCreate(BaseModel):
    product_code: str
    description: str = ""
    category: str = ""
    size: str = ""
    height_cm: float = 0
    depth_cm: float = 0
    width_cm: float = 0
    cbm: float = 0
    fob_price_usd: float = 0
    fob_price_gbp: float = 0
    warehouse_price_1: float = 0
    warehouse_price_2: float = 0
    image: str = ""
    images: List[str] = []

class ProductUpdate(BaseModel):
    product_code: Optional[str] = None
    description: Optional[str] = None
    category: Optional[str] = None
    size: Optional[str] = None
    height_cm: Optional[float] = None
    depth_cm: Optional[float] = None
    width_cm: Optional[float] = None
    cbm: Optional[float] = None
    fob_price_usd: Optional[float] = None
    fob_price_gbp: Optional[float] = None
    warehouse_price_1: Optional[float] = None
    warehouse_price_2: Optional[float] = None
    image: Optional[str] = None
    images: Optional[List[str]] = None

# Quotation Models
class QuotationItem(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    product_id: str = ""
    product_code: str = ""
    description: str = ""
    height_cm: float = 0
    depth_cm: float = 0
    width_cm: float = 0
    cbm: float = 0
    quantity: int = 1
    fob_price: float = 0
    total: float = 0

class Quotation(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    reference: str = ""
    customer_name: str = ""
    customer_email: str = ""
    date: str = ""
    currency: str = "USD"
    notes: str = ""
    items: List[QuotationItem] = []
    total_items: int = 0
    total_cbm: float = 0
    total_value: float = 0
    status: str = "draft"  # draft, sent, accepted, rejected
    created_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())
    updated_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())

class QuotationCreate(BaseModel):
    reference: str = ""
    customer_name: str = ""
    customer_email: str = ""
    date: str = ""
    currency: str = "USD"
    notes: str = ""
    items: List[dict] = []
    total_items: int = 0
    total_cbm: float = 0
    total_value: float = 0
    status: str = "draft"

class QuotationUpdate(BaseModel):
    reference: Optional[str] = None
    customer_name: Optional[str] = None
    customer_email: Optional[str] = None
    date: Optional[str] = None
    currency: Optional[str] = None
    notes: Optional[str] = None
    items: Optional[List[dict]] = None
    total_items: Optional[int] = None
    total_cbm: Optional[float] = None
    total_value: Optional[float] = None
    status: Optional[str] = None

# ============ ROUTES ============

@api_router.get("/")
async def root():
    return {"message": "JAIPUR Production Sheet API"}

# --- ORDERS ---

@api_router.get("/orders", response_model=List[Order])
async def get_orders():
    orders = await db.orders.find({}, {"_id": 0}).to_list(1000)
    return orders

@api_router.get("/orders/{order_id}", response_model=Order)
async def get_order(order_id: str):
    order = await db.orders.find_one({"id": order_id}, {"_id": 0})
    if not order:
        raise HTTPException(status_code=404, detail="Order not found")
    return order

@api_router.post("/orders", response_model=Order)
async def create_order(order_data: OrderCreate):
    order = Order(**order_data.model_dump())
    doc = order.model_dump()
    await db.orders.insert_one(doc)
    return order

@api_router.put("/orders/{order_id}", response_model=Order)
async def update_order(order_id: str, order_data: OrderUpdate):
    existing = await db.orders.find_one({"id": order_id}, {"_id": 0})
    if not existing:
        raise HTTPException(status_code=404, detail="Order not found")
    
    update_data = {k: v for k, v in order_data.model_dump().items() if v is not None}
    update_data["updated_at"] = datetime.now(timezone.utc).isoformat()
    
    if "items" in update_data:
        update_data["items"] = [item.model_dump() if hasattr(item, 'model_dump') else item for item in update_data["items"]]
    
    await db.orders.update_one({"id": order_id}, {"$set": update_data})
    updated = await db.orders.find_one({"id": order_id}, {"_id": 0})
    return updated

@api_router.delete("/orders/{order_id}")
async def delete_order(order_id: str):
    result = await db.orders.delete_one({"id": order_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Order not found")
    return {"message": "Order deleted"}

# --- LEATHER LIBRARY ---

@api_router.get("/leather-library", response_model=List[LeatherLibraryItem])
async def get_leather_library():
    items = await db.leather_library.find({}, {"_id": 0}).to_list(1000)
    return items

@api_router.post("/leather-library", response_model=LeatherLibraryItem)
async def create_leather_item(item: LeatherLibraryItem):
    doc = item.model_dump()
    await db.leather_library.insert_one(doc)
    return item

@api_router.put("/leather-library/{item_id}", response_model=LeatherLibraryItem)
async def update_leather_item(item_id: str, item: LeatherLibraryItem):
    await db.leather_library.update_one({"id": item_id}, {"$set": item.model_dump()})
    return item

@api_router.delete("/leather-library/{item_id}")
async def delete_leather_item(item_id: str):
    result = await db.leather_library.delete_one({"id": item_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Item not found")
    return {"message": "Item deleted"}

@api_router.post("/leather-library/upload-excel")
async def upload_leather_excel(file: UploadFile = File(...)):
    """Upload leather items from Excel file"""
    if not file.filename.endswith(('.xlsx', '.xls')):
        raise HTTPException(status_code=400, detail="File must be an Excel file")
    
    try:
        contents = await file.read()
        df = pd.read_excel(io.BytesIO(contents))
        df.columns = df.columns.str.lower().str.strip()
        
        column_mapping = {
            'code': 'code', 'leather code': 'code', 'item code': 'code',
            'name': 'name', 'leather name': 'name',
            'description': 'description', 'desc': 'description',
            'color': 'color',
            'image': 'image_url', 'image link': 'image_url', 'photo': 'image_url', 'photo link': 'image_url'
        }
        df = df.rename(columns=column_mapping)
        
        created = []
        for idx, row in df.iterrows():
            code = str(row.get('code', '')).strip()
            if not code or code == 'nan':
                continue
                
            item_data = {
                'id': str(uuid.uuid4()),
                'code': code.upper(),
                'name': str(row.get('name', '')).strip() if pd.notna(row.get('name')) else '',
                'description': str(row.get('description', '')).strip() if pd.notna(row.get('description')) else '',
                'color': str(row.get('color', '')).strip() if pd.notna(row.get('color')) else '',
                'image': '',
                'created_at': datetime.now(timezone.utc).isoformat()
            }
            
            # Fetch image from URL if provided
            image_url = str(row.get('image_url', '')).strip() if pd.notna(row.get('image_url')) else ''
            if image_url and image_url.startswith('http'):
                try:
                    async with httpx.AsyncClient(timeout=10.0) as client:
                        response = await client.get(image_url)
                        if response.status_code == 200 and 'image' in response.headers.get('content-type', ''):
                            item_data['image'] = f"data:{response.headers['content-type']};base64,{base64.b64encode(response.content).decode()}"
                except:
                    pass
            
            await db.leather_library.insert_one(item_data)
            created.append({'code': item_data['code'], 'name': item_data['name']})
        
        return {"message": f"{len(created)} items imported", "created": len(created), "items": created[:20]}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# --- FINISH LIBRARY ---

@api_router.get("/finish-library", response_model=List[FinishLibraryItem])
async def get_finish_library():
    items = await db.finish_library.find({}, {"_id": 0}).to_list(1000)
    return items

@api_router.post("/finish-library", response_model=FinishLibraryItem)
async def create_finish_item(item: FinishLibraryItem):
    doc = item.model_dump()
    await db.finish_library.insert_one(doc)
    return item

@api_router.put("/finish-library/{item_id}", response_model=FinishLibraryItem)
async def update_finish_item(item_id: str, item: FinishLibraryItem):
    await db.finish_library.update_one({"id": item_id}, {"$set": item.model_dump()})
    return item

@api_router.delete("/finish-library/{item_id}")
async def delete_finish_item(item_id: str):
    result = await db.finish_library.delete_one({"id": item_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Item not found")
    return {"message": "Item deleted"}

@api_router.post("/finish-library/upload-excel")
async def upload_finish_excel(file: UploadFile = File(...)):
    """Upload finish items from Excel file"""
    if not file.filename.endswith(('.xlsx', '.xls')):
        raise HTTPException(status_code=400, detail="File must be an Excel file")
    
    try:
        contents = await file.read()
        df = pd.read_excel(io.BytesIO(contents))
        df.columns = df.columns.str.lower().str.strip()
        
        column_mapping = {
            'code': 'code', 'finish code': 'code', 'item code': 'code',
            'name': 'name', 'finish name': 'name',
            'description': 'description', 'desc': 'description',
            'color': 'color',
            'image': 'image_url', 'image link': 'image_url', 'photo': 'image_url', 'photo link': 'image_url'
        }
        df = df.rename(columns=column_mapping)
        
        created = []
        for idx, row in df.iterrows():
            code = str(row.get('code', '')).strip()
            if not code or code == 'nan':
                continue
                
            item_data = {
                'id': str(uuid.uuid4()),
                'code': code.upper(),
                'name': str(row.get('name', '')).strip() if pd.notna(row.get('name')) else '',
                'description': str(row.get('description', '')).strip() if pd.notna(row.get('description')) else '',
                'color': str(row.get('color', '')).strip() if pd.notna(row.get('color')) else '',
                'image': '',
                'created_at': datetime.now(timezone.utc).isoformat()
            }
            
            # Fetch image from URL if provided
            image_url = str(row.get('image_url', '')).strip() if pd.notna(row.get('image_url')) else ''
            if image_url and image_url.startswith('http'):
                try:
                    async with httpx.AsyncClient(timeout=10.0) as client:
                        response = await client.get(image_url)
                        if response.status_code == 200 and 'image' in response.headers.get('content-type', ''):
                            item_data['image'] = f"data:{response.headers['content-type']};base64,{base64.b64encode(response.content).decode()}"
                except:
                    pass
            
            await db.finish_library.insert_one(item_data)
            created.append({'code': item_data['code'], 'name': item_data['name']})
        
        return {"message": f"{len(created)} items imported", "created": len(created), "items": created[:20]}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# --- TEMPLATE SETTINGS ---

@api_router.get("/template-settings", response_model=TemplateSettings)
async def get_template_settings():
    settings = await db.template_settings.find_one({"id": "default"}, {"_id": 0})
    if not settings:
        default_settings = TemplateSettings()
        await db.template_settings.insert_one(default_settings.model_dump())
        return default_settings
    return settings

@api_router.put("/template-settings", response_model=TemplateSettings)
async def update_template_settings(settings: TemplateSettings):
    settings.id = "default"
    await db.template_settings.update_one(
        {"id": "default"}, 
        {"$set": settings.model_dump()}, 
        upsert=True
    )
    return settings

# --- FACTORIES ---

@api_router.get("/factories")
async def get_factories():
    factories = await db.factories.find({}, {"_id": 0}).to_list(100)
    if not factories:
        # Initialize with default factories
        default_factories = [
            {"id": "sae", "code": "SAE", "name": "Shekhawati Art Exports"},
            {"id": "cac", "code": "CAC", "name": "Country Art & Crafts"},
            {"id": "gae", "code": "GAE", "name": "Global Art Exports"},
        ]
        for factory in default_factories:
            await db.factories.insert_one(factory)
        return default_factories
    return factories

@api_router.post("/factories")
async def create_factory(factory: dict):
    factory_id = str(uuid.uuid4())
    factory_doc = {
        "id": factory_id,
        "code": factory.get("code", ""),
        "name": factory.get("name", ""),
    }
    # Check if factory already exists
    existing = await db.factories.find_one({"name": factory_doc["name"]}, {"_id": 0})
    if existing:
        return existing
    await db.factories.insert_one(factory_doc)
    # Return without _id
    return {"id": factory_doc["id"], "code": factory_doc["code"], "name": factory_doc["name"]}

@api_router.delete("/factories/{factory_id}")
async def delete_factory(factory_id: str):
    result = await db.factories.delete_one({"id": factory_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Factory not found")
    return {"message": "Factory deleted"}

@api_router.post("/factories/upload-excel")
async def upload_factories_excel(file: UploadFile = File(...)):
    """Upload factories from Excel file"""
    if not file.filename.endswith(('.xlsx', '.xls')):
        raise HTTPException(status_code=400, detail="File must be an Excel file")
    
    try:
        contents = await file.read()
        df = pd.read_excel(io.BytesIO(contents))
        df.columns = df.columns.str.lower().str.strip()
        
        column_mapping = {
            'code': 'code', 'factory code': 'code',
            'name': 'name', 'factory name': 'name',
        }
        df = df.rename(columns=column_mapping)
        
        created = []
        for idx, row in df.iterrows():
            code = str(row.get('code', '')).strip()
            if not code or code == 'nan':
                continue
                
            factory_doc = {
                'id': str(uuid.uuid4()),
                'code': code.upper(),
                'name': str(row.get('name', '')).strip() if pd.notna(row.get('name')) else '',
            }
            
            await db.factories.insert_one(factory_doc)
            created.append({'code': factory_doc['code'], 'name': factory_doc['name']})
        
        return {"message": f"{len(created)} factories imported", "created": len(created), "items": created[:20]}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# --- CATEGORIES ---

@api_router.get("/categories")
async def get_categories():
    """Get all categories from database, with default ones if empty"""
    categories = await db.categories.find({}, {"_id": 0}).to_list(100)
    if not categories:
        # Initialize with default categories
        default_categories = [
            {"id": "chair", "name": "Chair"},
            {"id": "sofa", "name": "Sofa"},
            {"id": "bar-chair", "name": "Bar Chair"},
            {"id": "table", "name": "Table"},
            {"id": "bed", "name": "Bed"},
            {"id": "cabinet", "name": "Cabinet"},
            {"id": "shelf", "name": "Shelf"},
            {"id": "other", "name": "Other"}
        ]
        for cat in default_categories:
            await db.categories.insert_one(cat)
        return default_categories
    return categories

@api_router.post("/categories")
async def create_category(category: dict):
    """Add a new category"""
    category_id = category.get("id") or str(uuid.uuid4())
    category_doc = {
        "id": category_id,
        "name": category.get("name", ""),
    }
    # Check if category already exists
    existing = await db.categories.find_one({"name": category_doc["name"]}, {"_id": 0})
    if existing:
        return existing
    await db.categories.insert_one(category_doc)
    # Return without _id
    return {"id": category_doc["id"], "name": category_doc["name"]}

@api_router.delete("/categories/{category_id}")
async def delete_category(category_id: str):
    """Delete a category"""
    result = await db.categories.delete_one({"id": category_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Category not found")
    return {"message": "Category deleted"}

# --- PDF EXPORT ---

def strip_html(text):
    """Remove HTML tags from text"""
    clean = re.compile('<.*?>')
    return re.sub(clean, '', text or '')

def format_date_ddmmyyyy(date_str):
    """Format date string to DD-MM-YYYY format"""
    if not date_str:
        return '-'
    try:
        # Try parsing ISO format first
        if 'T' in str(date_str):
            date_str = str(date_str).split('T')[0]
        parts = str(date_str).split('-')
        if len(parts) == 3:
            year, month, day = parts
            return f"{day}-{month}-{year}"
        return date_str
    except:
        return date_str

# JAIPUR Logo URL
JAIPUR_LOGO_URL = "https://customer-assets.emergentagent.com/job_furnipdf-maker/artifacts/mdh71t2g_WhatsApp%20Image%202025-12-22%20at%202.24.36%20PM.jpeg"

async def fetch_image_bytes(url: str) -> bytes:
    """Fetch image from URL and return bytes"""
    try:
        async with httpx.AsyncClient(timeout=10.0) as client:
            response = await client.get(url)
            if response.status_code == 200:
                return response.content
    except:
        pass
    return None

def generate_pdf(order: dict, settings: dict, logo_bytes: bytes = None) -> bytes:
    """Generate PDF that matches the Preview page layout exactly - with LARGE product images"""
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=A4)
    width, height = A4
    margin = settings.get('page_margin_mm', 12) * mm
    
    primary_color = HexColor(settings.get('primary_color', '#3d2c1e'))
    
    for idx, item in enumerate(order.get("items", [])):
        # === HEADER SECTION (Logo LEFT, Info table RIGHT) ===
        header_top = height - margin
        
        # Logo on LEFT - larger size
        logo_width = 80
        logo_height = 50
        if logo_bytes:
            try:
                from reportlab.lib.utils import ImageReader
                logo_img = ImageReader(io.BytesIO(logo_bytes))
                c.drawImage(logo_img, margin, header_top - logo_height, width=logo_width, height=logo_height, preserveAspectRatio=True)
            except:
                c.setFillColor(primary_color)
                c.setFont("Helvetica-Bold", 24)
                c.drawString(margin, header_top - 25, "JAIPUR")
                c.setFont("Helvetica-Oblique", 7)
                c.setFillColor(HexColor('#666666'))
                c.drawString(margin, header_top - 35, "A fine wood furniture company")
        else:
            c.setFillColor(primary_color)
            c.setFont("Helvetica-Bold", 24)
            c.drawString(margin, header_top - 25, "JAIPUR")
            c.setFont("Helvetica-Oblique", 7)
            c.setFillColor(HexColor('#666666'))
            c.drawString(margin, header_top - 35, "A fine wood furniture company")
        
        # Info table on RIGHT
        table_width = 140
        right_x = width - margin - table_width
        y = header_top - 5
        
        dates = [
            ("ENTRY DATE", format_date_ddmmyyyy(order.get('entry_date', 'N/A'))),
            ("FACTORY INFORM", format_date_ddmmyyyy(order.get('factory_inform_date', order.get('entry_date', 'N/A')))),
            ("FACTORY", order.get('factory', 'N/A')),
            ("SALES REF", order.get('sales_order_ref', 'N/A')),
            ("BUYER PO", order.get('buyer_po_ref', 'N/A')),
        ]
        
        row_height = 12
        table_height = len(dates) * row_height
        
        # Draw table border
        c.setStrokeColor(primary_color)
        c.setLineWidth(1)
        c.rect(right_x, y - table_height, table_width, table_height)
        
        # Draw table rows
        label_width = 55
        for i, (label, value) in enumerate(dates):
            row_y = y - (i + 1) * row_height
            
            # Label cell with background
            c.setFillColor(HexColor('#f5f0eb'))
            c.rect(right_x, row_y, label_width, row_height, fill=True, stroke=False)
            
            # Cell borders
            c.setStrokeColor(primary_color)
            c.rect(right_x, row_y, label_width, row_height, fill=False, stroke=True)
            c.rect(right_x + label_width, row_y, table_width - label_width, row_height, fill=False, stroke=True)
            
            # Label text - INCREASED FONT SIZE
            c.setFillColor(primary_color)
            c.setFont("Helvetica-Bold", 8)
            c.drawString(right_x + 2, row_y + 3, label)
            
            # Value text - INCREASED FONT SIZE
            c.setFillColor(HexColor('#333333'))
            c.setFont("Helvetica", 8)
            value_str = str(value)[:20] if value else '-'
            c.drawString(right_x + label_width + 2, row_y + 3, value_str)
        
        # Separator line under header
        separator_y = header_top - 55
        c.setStrokeColor(primary_color)
        c.setLineWidth(2)
        c.line(margin, separator_y, width - margin, separator_y)
        
        # === LARGE PRODUCT IMAGE SECTION (75% width) + Material Swatches (25%) ===
        content_y = separator_y - 8
        content_width = width - 2*margin
        img_section_width = content_width * 0.73
        material_section_width = content_width * 0.25
        material_x = margin + img_section_width + 8
        
        # LARGE Product Image Area - additional 30% height increase (374 * 1.3 = 486)
        img_height = 486  # Much larger image area
        c.setStrokeColor(HexColor('#dddddd'))
        c.setLineWidth(0.5)
        c.rect(margin, content_y - img_height, img_section_width, img_height)
        
        # Try to draw product image - LARGE
        product_image = item.get('product_image') or (item.get('images', [None])[0] if item.get('images') else None)
        if product_image:
            try:
                if product_image.startswith('data:image'):
                    img_data = product_image.split(',')[1]
                    img_bytes_data = base64.b64decode(img_data)
                    from reportlab.lib.utils import ImageReader
                    img = ImageReader(io.BytesIO(img_bytes_data))
                    # Draw image with padding
                    c.drawImage(img, margin + 5, content_y - img_height + 5, 
                               width=img_section_width - 10, height=img_height - 10, 
                               preserveAspectRatio=True)
            except Exception as e:
                c.setFillColor(HexColor('#888888'))
                c.setFont("Helvetica", 12)
                c.drawCentredString(margin + img_section_width/2, content_y - img_height/2, "Image Error")
        else:
            c.setFillColor(HexColor('#888888'))
            c.setFont("Helvetica", 12)
            c.drawCentredString(margin + img_section_width/2, content_y - img_height/2, "No Image Available")
        
        # Material Swatches (25% width) - taller to match image height
        c.setStrokeColor(HexColor('#dddddd'))
        c.rect(material_x, content_y - img_height, material_section_width, img_height)
        
        # Leather swatch - LARGER (increased proportionally - 30% more)
        swatch_y = content_y - 10
        swatch_height = 156  # Larger swatches to match image height
        if item.get('leather_code') or item.get('leather_image'):
            if item.get('leather_image'):
                try:
                    if item['leather_image'].startswith('data:image'):
                        img_data = item['leather_image'].split(',')[1]
                        img_bytes_data = base64.b64decode(img_data)
                        from reportlab.lib.utils import ImageReader
                        img = ImageReader(io.BytesIO(img_bytes_data))
                        c.drawImage(img, material_x + 8, swatch_y - swatch_height, 
                                   width=material_section_width - 16, height=swatch_height - 5, 
                                   preserveAspectRatio=True)
                except:
                    c.setFillColor(HexColor('#8B4513'))
                    c.rect(material_x + 8, swatch_y - swatch_height, material_section_width - 16, swatch_height - 5, fill=True)
            else:
                c.setFillColor(HexColor('#8B4513'))
                c.rect(material_x + 8, swatch_y - swatch_height, material_section_width - 16, swatch_height - 5, fill=True)
            
            c.setFillColor(HexColor('#666666'))
            c.setFont("Helvetica", 6)
            c.drawCentredString(material_x + material_section_width/2, swatch_y - swatch_height - 8, "LEATHER")
            c.setFont("Helvetica-Bold", 7)
            c.setFillColor(HexColor('#333333'))
            c.drawCentredString(material_x + material_section_width/2, swatch_y - swatch_height - 18, item.get('leather_code', '-'))
            swatch_y -= swatch_height + 30
        
        # Finish swatch - LARGER
        if item.get('finish_code') or item.get('finish_image'):
            if item.get('finish_image'):
                try:
                    if item['finish_image'].startswith('data:image'):
                        img_data = item['finish_image'].split(',')[1]
                        img_bytes_data = base64.b64decode(img_data)
                        from reportlab.lib.utils import ImageReader
                        img = ImageReader(io.BytesIO(img_bytes_data))
                        c.drawImage(img, material_x + 8, swatch_y - swatch_height, 
                                   width=material_section_width - 16, height=swatch_height - 5, 
                                   preserveAspectRatio=True)
                except:
                    c.setFillColor(HexColor('#D4A574'))
                    c.rect(material_x + 8, swatch_y - swatch_height, material_section_width - 16, swatch_height - 5, fill=True)
            else:
                c.setFillColor(HexColor('#D4A574'))
                c.rect(material_x + 8, swatch_y - swatch_height, material_section_width - 16, swatch_height - 5, fill=True)
            
            c.setFillColor(HexColor('#666666'))
            c.setFont("Helvetica", 6)
            c.drawCentredString(material_x + material_section_width/2, swatch_y - swatch_height - 8, "FINISH")
            c.setFont("Helvetica-Bold", 7)
            c.setFillColor(HexColor('#333333'))
            c.drawCentredString(material_x + material_section_width/2, swatch_y - swatch_height - 18, item.get('finish_code', '-'))
        
        # No materials placeholder
        if not item.get('leather_code') and not item.get('leather_image') and not item.get('finish_code') and not item.get('finish_image'):
            c.setFillColor(HexColor('#888888'))
            c.setFont("Helvetica", 8)
            c.drawCentredString(material_x + material_section_width/2, content_y - img_height/2, "No material")
            c.drawCentredString(material_x + material_section_width/2, content_y - img_height/2 - 12, "swatches")
        
        # === NOTES SECTION (100% width) ===
        notes_y = content_y - img_height - 10
        notes_height = 60  # Increased height for better visibility
        c.setStrokeColor(primary_color)
        c.setLineWidth(1)
        c.rect(margin, notes_y - notes_height, content_width, notes_height)
        
        # Notes header
        c.setFillColor(primary_color)
        c.rect(margin, notes_y - 16, content_width, 16, fill=True)
        c.setFillColor(HexColor('#ffffff'))
        c.setFont("Helvetica-Bold", 10)  # Increased font
        c.drawString(margin + 5, notes_y - 12, "Notes:")
        
        # Notes content - INCREASED FONT SIZE
        c.setFillColor(HexColor('#333333'))
        c.setFont("Helvetica", 10)  # Increased from 8
        notes_text = strip_html(item.get('notes', ''))
        if notes_text:
            words = notes_text.split()
            line = ""
            line_y = notes_y - 30
            max_width = content_width - 20
            for word in words:
                test_line = line + " " + word if line else word
                if c.stringWidth(test_line, "Helvetica", 10) < max_width:
                    line = test_line
                else:
                    c.drawString(margin + 8, line_y, line)
                    line_y -= 13
                    line = word
                    if line_y < notes_y - notes_height + 5:
                        break
            if line and line_y >= notes_y - notes_height + 5:
                c.drawString(margin + 8, line_y, line)
        else:
            default_notes = []
            if item.get('category'): default_notes.append(f"Category: {item['category']}")
            if item.get('leather_code'): default_notes.append(f"Leather: {item['leather_code']}")
            if item.get('finish_code'): default_notes.append(f"Finish: {item['finish_code']}")
            if item.get('color_notes'): default_notes.append(f"Color: {item['color_notes']}")
            if item.get('wood_finish'): default_notes.append(f"Wood Finish: {item['wood_finish']}")
            line_y = notes_y - 30
            for note in default_notes[:3]:
                c.drawString(margin + 8, line_y, f"• {note}")
                line_y -= 13
        
        # === DETAILS TABLE (Bottom) ===
        table_y = notes_y - notes_height - 8
        
        # Table header with SIZE sub-headers - INCREASED HEIGHT AND FONT
        header_height = 28
        c.setFillColor(primary_color)
        c.rect(margin, table_y - header_height, content_width, header_height, fill=True)
        
        c.setFillColor(HexColor('#ffffff'))
        c.setFont("Helvetica-Bold", 10)  # Increased from 7
        
        # Column positions for 7 columns
        col_widths = [80, 170, 35, 35, 35, 50, 50]
        cols = [margin + 3]
        for w in col_widths[:-1]:
            cols.append(cols[-1] + w)
        
        # Draw headers
        headers = ["ITEM CODE", "DESCRIPTION", "H (cm)", "D (cm)", "W (cm)", "CBM", "Qty"]
        for i, header in enumerate(headers):
            c.drawString(cols[i], table_y - header_height + 10, header)
        
        # Table row
        cbm = item.get('cbm', 0)
        if item.get('cbm_auto', True):
            h = item.get('height_cm', 0) or 0
            d = item.get('depth_cm', 0) or 0
            w = item.get('width_cm', 0) or 0
            cbm = round((h * d * w) / 1000000, 4)
        
        row_y = table_y - header_height - 22  # Increased row height
        c.setStrokeColor(primary_color)
        c.rect(margin, row_y, content_width, 22)
        
        c.setFillColor(HexColor('#333333'))
        c.setFont("Courier-Bold", 11)  # Increased from 8
        c.drawString(cols[0], row_y + 7, str(item.get('product_code', '-')))
        
        c.setFont("Helvetica", 10)  # Increased from 8
        desc = item.get('description', '-')
        if item.get('color_notes'):
            desc = f"{desc} ({item['color_notes']})"
        if len(desc) > 45:
            desc = desc[:42] + "..."
        c.drawString(cols[1], row_y + 7, desc)
        c.drawString(cols[2], row_y + 7, str(item.get('height_cm', 0)))
        c.drawString(cols[3], row_y + 7, str(item.get('depth_cm', 0)))
        c.drawString(cols[4], row_y + 7, str(item.get('width_cm', 0)))
        c.drawString(cols[5], row_y + 7, str(cbm))
        c.setFont("Helvetica-Bold", 11)  # Increased from 8
        c.drawString(cols[6], row_y + 7, f"{item.get('quantity', 1)} Pcs")
        
        # Footer - INCREASED FONT
        c.setFillColor(HexColor('#666666'))
        c.setFont("Helvetica", 10)  # Increased from 8
        footer_text = f"Buyer: {order.get('buyer_name', 'N/A')} • PO: {order.get('buyer_po_ref', 'N/A')}"
        c.drawString(margin, margin + 10, footer_text)
        c.drawRightString(width - margin, margin + 10, f"Page {idx + 1} of {len(order.get('items', []))}")
        
        c.showPage()
    
    c.save()
    buffer.seek(0)
    return buffer.getvalue()

@api_router.get("/orders/{order_id}/export/pdf")
async def export_order_pdf(order_id: str):
    order = await db.orders.find_one({"id": order_id}, {"_id": 0})
    if not order:
        raise HTTPException(status_code=404, detail="Order not found")
    
    settings = await db.template_settings.find_one({"id": "default"}, {"_id": 0})
    if not settings:
        settings = TemplateSettings().model_dump()
    
    # Fetch logo image
    logo_bytes = await fetch_image_bytes(JAIPUR_LOGO_URL)
    
    pdf_bytes = generate_pdf(order, settings, logo_bytes)
    
    export_record = ExportRecord(
        order_id=order_id,
        export_type="pdf",
        filename=f"order_{order.get('sales_order_ref', order_id)}.pdf"
    )
    await db.exports.insert_one(export_record.model_dump())
    
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f"inline; filename=order_{order.get('sales_order_ref', order_id)}.pdf"}
    )

@api_router.get("/orders/{order_id}/preview-html")
async def preview_order_html(order_id: str):
    order = await db.orders.find_one({"id": order_id}, {"_id": 0})
    if not order:
        raise HTTPException(status_code=404, detail="Order not found")
    return {"html": "", "order": order}

# --- PPT EXPORT ---

@api_router.get("/orders/{order_id}/export/ppt")
async def export_order_ppt(order_id: str):
    order = await db.orders.find_one({"id": order_id}, {"_id": 0})
    if not order:
        raise HTTPException(status_code=404, detail="Order not found")
    
    settings = await db.template_settings.find_one({"id": "default"}, {"_id": 0})
    if not settings:
        settings = TemplateSettings().model_dump()
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Add logo image to title slide
    try:
        logo_bytes = await fetch_image_bytes(JAIPUR_LOGO_URL)
        if logo_bytes:
            logo_stream = io.BytesIO(logo_bytes)
            slide.shapes.add_picture(logo_stream, Inches(3.5), Inches(2), width=Inches(3))
    except:
        pass
    
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Production Sheet - {order.get('sales_order_ref', '')}"
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Buyer: {order.get('buyer_name', 'N/A')} | Date: {order.get('entry_date', 'N/A')}"
    p.font.size = Pt(18)
    
    # Create slide for each item
    for item in order.get("items", []):
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Header with product code
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(6), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{item.get('product_code', 'N/A')} - {item.get('description', '')}"
        p.font.size = Pt(20)
        p.font.bold = True
        
        # Info table on right
        info_text = f"""Entry Date: {order.get('entry_date', 'N/A')}
Factory: {order.get('factory', 'N/A')}
Sales Ref: {order.get('sales_order_ref', 'N/A')}
Buyer PO: {order.get('buyer_po_ref', 'N/A')}"""
        
        txBox = slide.shapes.add_textbox(Inches(7), Inches(0.2), Inches(2.8), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = info_text
        p.font.size = Pt(10)
        
        # Product image (left side, 70%)
        product_image = item.get('product_image') or (item.get('images', [None])[0] if item.get('images') else None)
        if product_image and product_image.startswith('data:image'):
            try:
                img_data = product_image.split(',')[1]
                img_bytes = base64.b64decode(img_data)
                img_stream = io.BytesIO(img_bytes)
                slide.shapes.add_picture(img_stream, Inches(0.3), Inches(1), width=Inches(6), height=Inches(4))
            except:
                # Add placeholder text
                txBox = slide.shapes.add_textbox(Inches(0.3), Inches(2.5), Inches(6), Inches(1))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = "No Image Available"
                p.font.size = Pt(14)
        else:
            txBox = slide.shapes.add_textbox(Inches(0.3), Inches(2.5), Inches(6), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "No Image Available"
            p.font.size = Pt(14)
        
        # Materials section (right side)
        materials_text = f"""MATERIALS:

Leather: {item.get('leather_code', 'N/A')}
Finish: {item.get('finish_code', 'N/A')}
Color Notes: {item.get('color_notes', 'N/A')}
Wood Finish: {item.get('wood_finish', 'N/A')}"""
        
        txBox = slide.shapes.add_textbox(Inches(6.5), Inches(1.2), Inches(3.3), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = materials_text
        p.font.size = Pt(11)
        
        # Notes section (below image)
        notes_plain = strip_html(item.get('notes', ''))
        if notes_plain:
            txBox = slide.shapes.add_textbox(Inches(0.3), Inches(5.2), Inches(9.4), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"Notes: {notes_plain}"
            p.font.size = Pt(10)
        
        # Details table at bottom
        cbm = item.get('cbm', 0)
        if item.get('cbm_auto', True):
            h = item.get('height_cm', 0) or 0
            d = item.get('depth_cm', 0) or 0
            w = item.get('width_cm', 0) or 0
            cbm = round((h * d * w) / 1000000, 4)
        
        table_text = f"""ITEM CODE: {item.get('product_code', '-')}  |  SIZE: {item.get('height_cm', 0)} × {item.get('depth_cm', 0)} × {item.get('width_cm', 0)} cm  |  CBM: {cbm}  |  QTY: {item.get('quantity', 1)} Pcs"""
        
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9.4), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = table_text
        p.font.size = Pt(11)
        p.font.bold = True
    
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    
    export_record = ExportRecord(
        order_id=order_id,
        export_type="ppt",
        filename=f"order_{order.get('sales_order_ref', order_id)}.pptx"
    )
    await db.exports.insert_one(export_record.model_dump())
    
    return StreamingResponse(
        ppt_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename=order_{order.get('sales_order_ref', order_id)}.pptx"}
    )

# --- EXPORT HISTORY ---

@api_router.get("/exports", response_model=List[ExportRecord])
async def get_exports():
    exports = await db.exports.find({}, {"_id": 0}).to_list(1000)
    return exports

@api_router.get("/exports/{order_id}", response_model=List[ExportRecord])
async def get_order_exports(order_id: str):
    exports = await db.exports.find({"order_id": order_id}, {"_id": 0}).to_list(100)
    return exports

# --- SAMPLE EXCEL TEMPLATES ---

@api_router.get("/templates/products-sample")
async def download_products_sample():
    """Download sample Excel template for products"""
    df = pd.DataFrame({
        'Product Code': ['IDR-8-180CM', 'KRL-2-180CM-LIGHT', 'BNC-4-90CM'],
        'Description': ['Induse Dining Table', 'Kerela Spider Leg Table', 'Branch Coffee Table'],
        'Size ( in Cm )': ['180*90 cm', '180 cm', '90*90 cm'],
        'H': [76, 76, 45],
        'D': [90, 90, 90],
        'W': [180, 180, 90],
        'CBM': [1.23, 1.23, 0.36],
        'FOB India Price $': [100, 120, 80],
        'FOB India Price £': [75, 90, 60],
        'Warehouse Price £700': [150, 180, 120],
        'Warehouse Price £2000': [180, 220, 150],
        'Photo Link': ['https://example.com/image1.jpg', 'https://example.com/image2.jpg', '']
    })
    
    buffer = io.BytesIO()
    df.to_excel(buffer, index=False, sheet_name='Products')
    buffer.seek(0)
    
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=products_sample.xlsx"}
    )

@api_router.get("/templates/leather-sample")
async def download_leather_sample():
    """Download sample Excel template for leather library"""
    df = pd.DataFrame({
        'Code': ['LTH-001', 'LTH-002', 'LTH-003'],
        'Name': ['Full Grain Tan', 'Nappa Black', 'Suede Brown'],
        'Description': ['Premium full grain leather', 'Soft nappa finish', 'Brushed suede texture'],
        'Color': ['#8B4513', '#000000', '#654321'],
        'Image Link': ['https://example.com/leather1.jpg', '', '']
    })
    
    buffer = io.BytesIO()
    df.to_excel(buffer, index=False, sheet_name='Leather')
    buffer.seek(0)
    
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=leather_sample.xlsx"}
    )

@api_router.get("/templates/finish-sample")
async def download_finish_sample():
    """Download sample Excel template for finish library"""
    df = pd.DataFrame({
        'Code': ['FIN-001', 'FIN-002', 'FIN-003'],
        'Name': ['Antique Brass', 'Matte Black', 'Polished Chrome'],
        'Description': ['Vintage brass finish', 'Modern matte finish', 'Shiny chrome plating'],
        'Color': ['#CD7F32', '#1C1C1C', '#C0C0C0'],
        'Image Link': ['https://example.com/finish1.jpg', '', '']
    })
    
    buffer = io.BytesIO()
    df.to_excel(buffer, index=False, sheet_name='Finish')
    buffer.seek(0)
    
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=finish_sample.xlsx"}
    )

@api_router.get("/templates/factories-sample")
async def download_factories_sample():
    """Download sample Excel template for factories"""
    df = pd.DataFrame({
        'Code': ['SAE', 'CAC', 'GAE', 'JFW'],
        'Name': ['Shekhawati Art Exports', 'Country Art & Crafts', 'Global Art Exports', 'Jaipur Fine Wood']
    })
    
    buffer = io.BytesIO()
    df.to_excel(buffer, index=False, sheet_name='Factories')
    buffer.seek(0)
    
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=factories_sample.xlsx"}
    )

# --- PRODUCTS ---

@api_router.get("/products", response_model=List[Product])
async def get_products():
    products = await db.products.find({}, {"_id": 0}).to_list(1000)
    return products

@api_router.get("/products/{product_id}", response_model=Product)
async def get_product(product_id: str):
    product = await db.products.find_one({"id": product_id}, {"_id": 0})
    if not product:
        raise HTTPException(status_code=404, detail="Product not found")
    return product

@api_router.post("/products", response_model=Product)
async def create_product(product_data: ProductCreate):
    product = Product(**product_data.model_dump())
    doc = product.model_dump()
    await db.products.insert_one(doc)
    return product

@api_router.put("/products/{product_id}", response_model=Product)
async def update_product(product_id: str, product_data: ProductUpdate):
    existing = await db.products.find_one({"id": product_id}, {"_id": 0})
    if not existing:
        raise HTTPException(status_code=404, detail="Product not found")
    
    update_data = {k: v for k, v in product_data.model_dump().items() if v is not None}
    update_data["updated_at"] = datetime.now(timezone.utc).isoformat()
    
    await db.products.update_one({"id": product_id}, {"$set": update_data})
    updated = await db.products.find_one({"id": product_id}, {"_id": 0})
    return updated

@api_router.delete("/products/{product_id}")
async def delete_product(product_id: str):
    result = await db.products.delete_one({"id": product_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Product not found")
    return {"message": "Product deleted"}

@api_router.post("/products/bulk")
async def bulk_create_products(products: List[ProductCreate]):
    created = []
    for product_data in products:
        product = Product(**product_data.model_dump())
        doc = product.model_dump()
        await db.products.insert_one(doc)
        created.append(product)
    return {"message": f"{len(created)} products created", "products": created}

@api_router.post("/products/upload-excel")
async def upload_products_excel(file: UploadFile = File(...)):
    """Upload products from Excel file with optional image URLs"""
    if not file.filename.endswith(('.xlsx', '.xls')):
        raise HTTPException(status_code=400, detail="File must be an Excel file (.xlsx or .xls)")
    
    try:
        # Read Excel file
        contents = await file.read()
        df = pd.read_excel(io.BytesIO(contents))
        
        # Check if first row contains actual headers (e.g., "Product Code", "Description")
        # This handles Excel files with merged header cells or wrong header detection
        first_row_values = df.iloc[0].astype(str).str.lower().tolist()
        header_keywords = ['product code', 'description', 'cbm', 'photo link', 'h', 'd', 'w']
        has_header_in_first_row = any(keyword in ' '.join(first_row_values) for keyword in header_keywords)
        
        if has_header_in_first_row or any('unnamed' in str(col).lower() for col in df.columns):
            # First row contains actual headers - use it and skip
            new_headers = df.iloc[0].astype(str).tolist()
            df = df.iloc[1:].reset_index(drop=True)
            df.columns = new_headers
        
        # Normalize column names (handle various formats)
        column_mapping = {
            'product code': 'product_code',
            'productcode': 'product_code',
            'item code': 'product_code',
            'itemcode': 'product_code',
            'code': 'product_code',
            'description': 'description',
            'desc': 'description',
            'size': 'size',
            'size ( in cm )': 'size',
            'size (cm)': 'size',
            'h': 'height_cm',
            'height': 'height_cm',
            'height_cm': 'height_cm',
            'd': 'depth_cm',
            'depth': 'depth_cm',
            'depth_cm': 'depth_cm',
            'w': 'width_cm',
            'width': 'width_cm',
            'width_cm': 'width_cm',
            'cbm': 'cbm',
            'fob india price $': 'fob_price_usd',
            'fob $': 'fob_price_usd',
            'fob_price_usd': 'fob_price_usd',
            'fob india price £': 'fob_price_gbp',
            'fob £': 'fob_price_gbp',
            'fob_price_gbp': 'fob_price_gbp',
            'warehouse price £700': 'warehouse_price_1',
            'warehouse_price_1': 'warehouse_price_1',
            'warehouse price £2000': 'warehouse_price_2',
            'warehouse_price_2': 'warehouse_price_2',
            'photo link': 'image_url',
            'photolink': 'image_url',
            'image': 'image_url',
            'image_url': 'image_url',
            'photo': 'image_url',
            'picture': 'image_url',
            'category': 'category',
        }
        
        # Rename columns
        df.columns = df.columns.str.lower().str.strip()
        df = df.rename(columns=column_mapping)
        
        created_products = []
        skipped = 0
        errors = []
        
        for idx, row in df.iterrows():
            try:
                # Skip rows without product code
                product_code = str(row.get('product_code', '')).strip()
                if not product_code or product_code == 'nan' or product_code == '':
                    skipped += 1
                    continue
                
                # Parse numeric fields safely
                def safe_float(val, default=0):
                    try:
                        if pd.isna(val) or val == '' or val == 'nan':
                            return default
                        return float(val)
                    except:
                        return default
                
                # Build product data
                product_data = {
                    'product_code': product_code.upper(),
                    'description': str(row.get('description', '')).strip() if pd.notna(row.get('description')) else '',
                    'size': str(row.get('size', '')).strip() if pd.notna(row.get('size')) else '',
                    'category': str(row.get('category', '')).strip() if pd.notna(row.get('category')) else '',
                    'height_cm': safe_float(row.get('height_cm')),
                    'depth_cm': safe_float(row.get('depth_cm')),
                    'width_cm': safe_float(row.get('width_cm')),
                    'cbm': safe_float(row.get('cbm')),
                    'fob_price_usd': safe_float(row.get('fob_price_usd')),
                    'fob_price_gbp': safe_float(row.get('fob_price_gbp')),
                    'warehouse_price_1': safe_float(row.get('warehouse_price_1')),
                    'warehouse_price_2': safe_float(row.get('warehouse_price_2')),
                    'image': '',
                    'images': []
                }
                
                # Handle image URL - try to fetch and convert to base64
                image_url = str(row.get('image_url', '')).strip() if pd.notna(row.get('image_url')) else ''
                if image_url and image_url != 'nan' and image_url != '#REF!' and image_url.startswith('http'):
                    try:
                        async with httpx.AsyncClient(timeout=10.0) as client:
                            response = await client.get(image_url)
                            if response.status_code == 200:
                                content_type = response.headers.get('content-type', 'image/jpeg')
                                if 'image' in content_type:
                                    image_base64 = base64.b64encode(response.content).decode('utf-8')
                                    product_data['image'] = f"data:{content_type};base64,{image_base64}"
                    except Exception as img_error:
                        # Log but continue without image
                        pass
                
                # Create product
                product = Product(**product_data)
                doc = product.model_dump()
                await db.products.insert_one(doc)
                created_products.append({
                    'product_code': product.product_code,
                    'description': product.description
                })
                
            except Exception as row_error:
                errors.append(f"Row {idx + 2}: {str(row_error)}")
        
        return {
            "message": f"Successfully imported {len(created_products)} products",
            "created": len(created_products),
            "skipped": skipped,
            "errors": errors[:10] if errors else [],  # Return first 10 errors only
            "products": created_products[:20]  # Return first 20 products
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing Excel file: {str(e)}")

# --- DASHBOARD STATS ---

@api_router.get("/dashboard/stats")
async def get_dashboard_stats():
    total_orders = await db.orders.count_documents({})
    draft_orders = await db.orders.count_documents({"status": "Draft"})
    in_production = await db.orders.count_documents({"status": "In Production"})
    completed = await db.orders.count_documents({"status": "Done"})
    
    recent_orders = await db.orders.find({}, {"_id": 0}).sort("created_at", -1).limit(5).to_list(5)
    
    return {
        "total_orders": total_orders,
        "draft_orders": draft_orders,
        "in_production": in_production,
        "completed": completed,
        "recent_orders": recent_orders
    }

# --- QUOTATIONS ---

@api_router.get("/quotations", response_model=List[Quotation])
async def get_quotations():
    quotations = await db.quotations.find({}, {"_id": 0}).sort("created_at", -1).to_list(1000)
    return quotations

@api_router.get("/quotations/{quotation_id}", response_model=Quotation)
async def get_quotation(quotation_id: str):
    quotation = await db.quotations.find_one({"id": quotation_id}, {"_id": 0})
    if not quotation:
        raise HTTPException(status_code=404, detail="Quotation not found")
    return quotation

@api_router.post("/quotations", response_model=Quotation)
async def create_quotation(quotation: QuotationCreate):
    doc = Quotation(**quotation.model_dump()).model_dump()
    await db.quotations.insert_one(doc)
    return doc

@api_router.put("/quotations/{quotation_id}", response_model=Quotation)
async def update_quotation(quotation_id: str, quotation: QuotationUpdate):
    existing = await db.quotations.find_one({"id": quotation_id}, {"_id": 0})
    if not existing:
        raise HTTPException(status_code=404, detail="Quotation not found")
    
    update_data = {k: v for k, v in quotation.model_dump().items() if v is not None}
    update_data["updated_at"] = datetime.now(timezone.utc).isoformat()
    
    await db.quotations.update_one({"id": quotation_id}, {"$set": update_data})
    updated = await db.quotations.find_one({"id": quotation_id}, {"_id": 0})
    return updated

@api_router.delete("/quotations/{quotation_id}")
async def delete_quotation(quotation_id: str):
    result = await db.quotations.delete_one({"id": quotation_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Quotation not found")
    return {"message": "Quotation deleted successfully"}

@api_router.post("/quotations/{quotation_id}/duplicate", response_model=Quotation)
async def duplicate_quotation(quotation_id: str):
    """Duplicate an existing quotation for reuse"""
    existing = await db.quotations.find_one({"id": quotation_id}, {"_id": 0})
    if not existing:
        raise HTTPException(status_code=404, detail="Quotation not found")
    
    # Create new quotation based on existing one
    new_quotation = {
        **existing,
        "id": str(uuid.uuid4()),
        "reference": f"{existing.get('reference', 'QT')}-COPY",
        "date": datetime.now(timezone.utc).strftime("%Y-%m-%d"),
        "status": "draft",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "updated_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.quotations.insert_one(new_quotation)
    return new_quotation

# Include the router in the main app
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
